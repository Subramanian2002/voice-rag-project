import pymupdf
from pptx import Presentation


# extract text from PDF file
def extract_pdf_text(file_path: str) -> str:
    document = pymupdf.open(file_path)

    text = ""

    for page in document:
        text += page.get_text()

    document.close()

    return text

# extract text from TXT file
def extract_txt_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()


# extract text from PPT file
def extract_pptx_text(file_path: str) -> str:
    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

